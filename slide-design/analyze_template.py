# -*- coding: utf-8 -*-
"""
会社FMTテンプレートの分析スクリプト
フォント、カラーパレット、レイアウト情報を抽出
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from collections import defaultdict
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

template_path = r"C:\Users\KazuhideAkitake\.claude\skills\slide-design\templates\company_format.pptx"

def analyze_template(prs):
    """テンプレートを詳細分析"""

    print("=" * 60)
    print("会社FMTテンプレート分析結果")
    print("=" * 60)

    # スライド数
    print(f"\n■ スライド数: {len(prs.slides)}")

    # スライドサイズ
    print(f"\n■ スライドサイズ:")
    print(f"  幅: {prs.slide_width.inches:.2f} inches")
    print(f"  高さ: {prs.slide_height.inches:.2f} inches")

    # フォント情報を収集
    fonts = defaultdict(int)
    font_sizes = defaultdict(int)
    colors = defaultdict(int)

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        font = run.font
                        if font.name:
                            fonts[font.name] += 1
                        if font.size:
                            font_sizes[f"{font.size.pt}pt"] += 1
                        try:
                            if font.color and font.color.rgb:
                                colors[str(font.color.rgb)] += 1
                        except:
                            pass

    # フォント
    print(f"\n■ 使用フォント（出現回数順）:")
    for font, count in sorted(fonts.items(), key=lambda x: -x[1])[:10]:
        print(f"  {font}: {count}回")

    # フォントサイズ
    print(f"\n■ 使用フォントサイズ（出現回数順）:")
    for size, count in sorted(font_sizes.items(), key=lambda x: -x[1])[:10]:
        print(f"  {size}: {count}回")

    # カラー
    print(f"\n■ 使用カラー（出現回数順）:")
    for color, count in sorted(colors.items(), key=lambda x: -x[1])[:15]:
        print(f"  #{color}: {count}回")

    # スライドマスタ情報
    print(f"\n■ スライドマスタ数: {len(prs.slide_masters)}")
    for i, master in enumerate(prs.slide_masters):
        print(f"\n  マスタ {i+1}:")
        print(f"    レイアウト数: {len(master.slide_layouts)}")
        for j, layout in enumerate(master.slide_layouts):
            print(f"      [{j}] {layout.name}")

    # 代表的なスライドの構成分析
    print("\n" + "=" * 60)
    print("代表スライドのレイアウト分析")
    print("=" * 60)

    # 最初の5スライドを分析
    for slide_idx in range(min(5, len(prs.slides))):
        slide = prs.slides[slide_idx]
        print(f"\n■ スライド {slide_idx + 1}: {slide.slide_layout.name}")

        shapes_info = []
        for shape in slide.shapes:
            info = {
                "name": shape.name,
                "left": shape.left.inches,
                "top": shape.top.inches,
                "width": shape.width.inches,
                "height": shape.height.inches,
            }

            if shape.has_text_frame:
                text = shape.text_frame.text[:50].replace('\n', ' ').strip()
                if text:
                    info["text"] = text

                    # フォント情報
                    for para in shape.text_frame.paragraphs:
                        if para.runs:
                            run = para.runs[0]
                            font = run.font
                            info["font_name"] = font.name if font.name else "inherit"
                            info["font_size"] = f"{font.size.pt}pt" if font.size else "inherit"
                            info["font_bold"] = font.bold
                            break

            shapes_info.append(info)

        # 上から順にソート
        shapes_info.sort(key=lambda x: (x["top"], x["left"]))

        for s in shapes_info:
            if "text" in s:
                print(f"  [{s['left']:.2f}, {s['top']:.2f}] w={s['width']:.2f} h={s['height']:.2f}")
                print(f"    Text: {s['text']}")
                if "font_name" in s:
                    print(f"    Font: {s.get('font_name', '?')} / {s.get('font_size', '?')} / Bold={s.get('font_bold', '?')}")

# テンプレートを読み込み
prs = Presentation(template_path)
analyze_template(prs)
